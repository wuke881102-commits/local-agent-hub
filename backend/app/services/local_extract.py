"""把本地文档（pdf / word / excel / ppt / html / csv）抽取成 Markdown 文本。

供「本地目录」的内容生产用：非图片文件先确定性抽取文字，再交给文本模型重组。
图片不在此处理（由 Agent 走视觉模型）。复用 office_reader / pdf_reader 的本地解析能力；
PPT 用 python-pptx（可选依赖，缺失则返回提示文案而非抛错）。
"""
from __future__ import annotations

from pathlib import Path

from . import office_reader, pdf_reader

_MAX_CHARS = 100000      # 单文件抽取上限，避免超大文档爆上下文
_MAX_SHEET_ROWS = 60     # excel 每表最多取的数据行（喂模型够用）


def extract_markdown(path: str | Path, kind: str) -> str:
    """按 kind 抽取本地文件正文为 Markdown；不可解析时返回简短提示串（不抛错）。"""
    p = Path(path)
    if not p.is_file():
        return ""
    try:
        if kind == "pdf":
            md = _pdf(p)
        elif kind == "word":
            md = _word(p)
        elif kind == "excel":
            md = _excel(p)
        elif kind == "ppt":
            md = _ppt(p)
        else:
            return ""
    except Exception as e:  # noqa: BLE001
        return f"（{p.name} 解析失败：{type(e).__name__}）"
    md = (md or "").strip()
    if not md:
        return f"（{p.name} 未抽取到文字内容，可能是扫描件/空文件）"
    if len(md) > _MAX_CHARS:
        md = md[:_MAX_CHARS].rstrip() + "\n\n…（内容较长，已截断）"
    return md


def _pdf(p: Path) -> str:
    if not pdf_reader.available():
        return "（缺少 PyMuPDF，无法解析 PDF）"
    # 只取文字层（max_vision_pages=0 → 不渲染、不调视觉），快速且无外部依赖。
    res = pdf_reader.extract(p, max_vision_pages=0, want_figures=False)
    return pdf_reader.assemble_full_text(res.get("pages") or [])


def _word(p: Path) -> str:
    ext = p.suffix.lower()
    if ext == ".docx":
        if not office_reader.available("word"):
            return "（缺少 python-docx，无法解析 Word）"
        return office_reader.parse_docx(p).get("markdown", "")
    return "（暂不支持旧版 .doc，请另存为 .docx）"


def _excel(p: Path) -> str:
    ext = p.suffix.lower()
    if ext == ".csv":
        try:
            return p.read_text(encoding="utf-8-sig", errors="replace")
        except OSError as e:
            return f"（CSV 读取失败：{e}）"
    if ext == ".xls":
        return "（暂不支持旧版 .xls，请另存为 .xlsx）"
    if not office_reader.available("excel"):
        return "（缺少 openpyxl，无法解析 Excel）"
    sheets = office_reader.parse_xlsx(p)
    parts: list[str] = []
    for s in sheets:
        headers = [str(h) for h in (s.get("headers") or [])]
        rows = (s.get("rows") or [])[:_MAX_SHEET_ROWS]
        if not headers:
            continue
        parts.append(f"## {s.get('name') or '工作表'}")
        parts.append("| " + " | ".join(headers) + " |")
        parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for r in rows:
            cells = ["" if c is None else str(c).replace("\n", " ") for c in r]
            cells += [""] * (len(headers) - len(cells))
            parts.append("| " + " | ".join(cells[:len(headers)]) + " |")
    return "\n".join(parts)


def _ppt(p: Path) -> str:
    if p.suffix.lower() != ".pptx":
        return "（暂不支持旧版 .ppt，请另存为 .pptx）"
    try:
        from pptx import Presentation  # noqa: PLC0415
    except Exception:  # noqa: BLE001
        return "（缺少 python-pptx，无法解析 PPT）"
    prs = Presentation(str(p))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    lines.append(shape.text_frame.text.strip())
                elif getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        lines.append("| " + " | ".join(cells) + " |")
            except Exception:  # noqa: BLE001
                continue
        if lines:
            parts.append(f"## 第 {i} 页\n" + "\n".join(lines))
    return "\n\n".join(parts)
